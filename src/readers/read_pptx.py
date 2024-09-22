from pptx import Presentation


def read_pptx(path):
    prs = Presentation(path)
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                fullText.append(paragraph.text)

    return "\n".join(fullText).replace("\n", " ")
